"""
Live Vibe Coding Club — Brand PPT Template Generator

ブランドデザインシステムに基づいたPowerPointテンプレートを生成する。

スライドレイアウト:
  0. タイトルスライド（ダーク背景 + グラデーションライン）
  1. セクション区切り
  2. コンテンツスライド（タイトル + 本文）
  3. 2カラムスライド
  4. コードブロックスライド
  5. 画像 + テキストスライド
  6. エンディングスライド
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ============================================================
# Brand Design Tokens
# ============================================================

BRAND = {
    "primary": "#8B5CF6",
    "primary_dark": "#6D28D9",
    "secondary": "#06B6D4",
    "accent": "#EC4899",
    "neutral_950": "#050507",
    "neutral_900": "#0A0A14",
    "neutral_800": "#12121F",
    "neutral_700": "#1E1E32",
    "neutral_600": "#33334A",
    "neutral_400": "#71718A",
    "neutral_300": "#A1A1B5",
    "neutral_200": "#D4D4DE",
    "neutral_100": "#EDEDF3",
    "neutral_50": "#F8F8FC",
    "success": "#10B981",
    "danger": "#EF4444",
    "text_white": "#F8F8FC",
    "text_muted": "#A1A1B5",
}

FONT_DISPLAY = "Outfit"
FONT_BODY = "Outfit"
FONT_MONO = "JetBrains Mono"
FONT_JP = "Noto Sans JP"

# Slide dimensions: 16:9
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_slide_bg(slide, hex_color: str):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(hex_color)


def add_gradient_line(slide, top, width=None, height=Pt(3)):
    """ブランドグラデーションを模したカラーバー（3色分割）を追加"""
    if width is None:
        width = SLIDE_WIDTH
    segment_w = int(width / 3)

    colors = [BRAND["primary"], BRAND["secondary"], BRAND["accent"]]
    for i, color in enumerate(colors):
        left = int(segment_w * i)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, segment_w, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(color)
        shape.line.fill.background()


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text="",
    font_name=FONT_DISPLAY,
    font_size=Pt(18),
    font_color=BRAND["text_white"],
    bold=False,
    alignment=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    if anchor:
        tf.paragraphs[0].alignment = alignment

    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = font_size
    p.font.color.rgb = hex_to_rgb(font_color)
    p.font.bold = bold
    p.alignment = alignment

    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color, corner_radius=Pt(12)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
    shape.line.fill.background()
    return shape


def add_placeholder_image(slide, left, top, width, height):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(BRAND["neutral_700"])
    shape.line.color.rgb = hex_to_rgb(BRAND["neutral_600"])
    shape.line.width = Pt(1)

    # 中央にテキスト
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "画像を挿入"
    p.font.name = FONT_JP
    p.font.size = Pt(14)
    p.font.color.rgb = hex_to_rgb(BRAND["neutral_400"])
    p.alignment = PP_ALIGN.CENTER

    return shape


# ============================================================
# Slide Generators
# ============================================================


def create_title_slide(prs: Presentation):
    """スライド0: タイトルスライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_bg(slide, BRAND["neutral_950"])

    # Subtle glow effect (large semi-transparent circle)
    glow = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(3), Inches(0.5),
        Inches(7), Inches(7),
    )
    glow.fill.solid()
    glow.fill.fore_color.rgb = hex_to_rgb(BRAND["primary"])
    glow.fill.fore_color.brightness = 0.85
    glow.line.fill.background()
    # Make it semi-transparent via alpha
    from pptx.oxml.ns import qn
    spPr = glow._element.find(qn("a:solidFill"))
    if spPr is None:
        # Access via shape properties
        fill_elem = glow._element.find(qn("p:spPr"))
        if fill_elem is not None:
            solid = fill_elem.find(qn("a:solidFill"))
            if solid is not None:
                srgb = solid.find(qn("a:srgbClr"))
                if srgb is not None:
                    alpha = srgb.makeelement(qn("a:alpha"), {"val": "8000"})
                    srgb.append(alpha)

    # Gradient bar at top
    add_gradient_line(slide, top=Pt(0), height=Pt(4))

    # Brand name (small, top-left)
    add_textbox(
        slide,
        left=Inches(1), top=Inches(1.2),
        width=Inches(6), height=Inches(0.6),
        text="LIVE VIBE CODING CLUB",
        font_name=FONT_DISPLAY, font_size=Pt(14),
        font_color=BRAND["neutral_400"],
        bold=True,
    )

    # Main title
    add_textbox(
        slide,
        left=Inches(1), top=Inches(2.2),
        width=Inches(10), height=Inches(2),
        text="プレゼンテーションタイトル",
        font_name=FONT_JP, font_size=Pt(48),
        font_color=BRAND["text_white"],
        bold=True,
    )

    # Subtitle
    add_textbox(
        slide,
        left=Inches(1), top=Inches(4.4),
        width=Inches(10), height=Inches(1),
        text="サブタイトルをここに入力",
        font_name=FONT_JP, font_size=Pt(22),
        font_color=BRAND["text_muted"],
    )

    # Date & presenter
    add_textbox(
        slide,
        left=Inches(1), top=Inches(6.0),
        width=Inches(6), height=Inches(0.5),
        text="2026.02.18  |  発表者名",
        font_name=FONT_BODY, font_size=Pt(14),
        font_color=BRAND["neutral_400"],
    )

    # Gradient bar at bottom
    add_gradient_line(slide, top=SLIDE_HEIGHT - Pt(4), height=Pt(4))


def create_section_slide(prs: Presentation):
    """スライド1: セクション区切り"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND["neutral_900"])

    # Left accent bar (primary color)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(2.5),
        Pt(6), Inches(2.5),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_to_rgb(BRAND["primary"])
    bar.line.fill.background()

    # Section number
    add_textbox(
        slide,
        left=Inches(1.2), top=Inches(2.3),
        width=Inches(3), height=Inches(0.7),
        text="01",
        font_name=FONT_DISPLAY, font_size=Pt(24),
        font_color=BRAND["primary"],
        bold=True,
    )

    # Section title
    add_textbox(
        slide,
        left=Inches(1.2), top=Inches(3.0),
        width=Inches(10), height=Inches(1.5),
        text="セクションタイトル",
        font_name=FONT_JP, font_size=Pt(40),
        font_color=BRAND["text_white"],
        bold=True,
    )

    # Section description
    add_textbox(
        slide,
        left=Inches(1.2), top=Inches(4.5),
        width=Inches(8), height=Inches(0.8),
        text="セクションの概要をここに記載します",
        font_name=FONT_JP, font_size=Pt(16),
        font_color=BRAND["text_muted"],
    )

    # Gradient bar at bottom
    add_gradient_line(slide, top=SLIDE_HEIGHT - Pt(4), height=Pt(4))


def create_content_slide(prs: Presentation):
    """スライド2: コンテンツスライド（タイトル + 本文）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND["neutral_950"])

    # Top gradient bar
    add_gradient_line(slide, top=Pt(0), height=Pt(3))

    # Slide title
    add_textbox(
        slide,
        left=Inches(0.8), top=Inches(0.6),
        width=Inches(11), height=Inches(0.9),
        text="スライドタイトル",
        font_name=FONT_JP, font_size=Pt(32),
        font_color=BRAND["text_white"],
        bold=True,
    )

    # Divider line
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(1.5),
        Inches(2), Pt(2),
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = hex_to_rgb(BRAND["primary"])
    divider.line.fill.background()

    # Body text area
    txBox = add_textbox(
        slide,
        left=Inches(0.8), top=Inches(1.9),
        width=Inches(11), height=Inches(4.8),
        text="",
        font_name=FONT_JP, font_size=Pt(18),
        font_color=BRAND["text_white"],
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    # Add sample body text with multiple paragraphs
    p = tf.paragraphs[0]
    p.text = "本文テキストをここに入力します。"
    p.font.name = FONT_JP
    p.font.size = Pt(18)
    p.font.color.rgb = hex_to_rgb(BRAND["text_white"])
    p.space_after = Pt(12)

    p2 = tf.add_paragraph()
    p2.text = "・箇条書き項目 1"
    p2.font.name = FONT_JP
    p2.font.size = Pt(16)
    p2.font.color.rgb = hex_to_rgb(BRAND["neutral_200"])
    p2.space_after = Pt(6)

    p3 = tf.add_paragraph()
    p3.text = "・箇条書き項目 2"
    p3.font.name = FONT_JP
    p3.font.size = Pt(16)
    p3.font.color.rgb = hex_to_rgb(BRAND["neutral_200"])
    p3.space_after = Pt(6)

    p4 = tf.add_paragraph()
    p4.text = "・箇条書き項目 3"
    p4.font.name = FONT_JP
    p4.font.size = Pt(16)
    p4.font.color.rgb = hex_to_rgb(BRAND["neutral_200"])

    # Page number area (bottom right)
    add_textbox(
        slide,
        left=Inches(11.5), top=Inches(6.8),
        width=Inches(1.5), height=Inches(0.4),
        text="03",
        font_name=FONT_DISPLAY, font_size=Pt(12),
        font_color=BRAND["neutral_400"],
        alignment=PP_ALIGN.RIGHT,
    )


def create_two_column_slide(prs: Presentation):
    """スライド3: 2カラムレイアウト"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND["neutral_950"])

    # Top gradient bar
    add_gradient_line(slide, top=Pt(0), height=Pt(3))

    # Slide title
    add_textbox(
        slide,
        left=Inches(0.8), top=Inches(0.6),
        width=Inches(11), height=Inches(0.9),
        text="2カラム レイアウト",
        font_name=FONT_JP, font_size=Pt(32),
        font_color=BRAND["text_white"],
        bold=True,
    )

    # Divider
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(1.5),
        Inches(2), Pt(2),
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = hex_to_rgb(BRAND["primary"])
    divider.line.fill.background()

    # Left column card
    left_card = add_rounded_rect(
        slide,
        Inches(0.8), Inches(2.0),
        Inches(5.6), Inches(4.5),
        BRAND["neutral_800"],
    )
    # Left column title
    add_textbox(
        slide,
        left=Inches(1.2), top=Inches(2.3),
        width=Inches(4.8), height=Inches(0.7),
        text="左カラム見出し",
        font_name=FONT_JP, font_size=Pt(22),
        font_color=BRAND["text_white"],
        bold=True,
    )
    # Left column body
    add_textbox(
        slide,
        left=Inches(1.2), top=Inches(3.1),
        width=Inches(4.8), height=Inches(3),
        text="左カラムの本文テキスト。ポイントや説明を記載します。",
        font_name=FONT_JP, font_size=Pt(16),
        font_color=BRAND["neutral_200"],
    )

    # Right column card
    right_card = add_rounded_rect(
        slide,
        Inches(6.8), Inches(2.0),
        Inches(5.6), Inches(4.5),
        BRAND["neutral_800"],
    )
    # Right column title
    add_textbox(
        slide,
        left=Inches(7.2), top=Inches(2.3),
        width=Inches(4.8), height=Inches(0.7),
        text="右カラム見出し",
        font_name=FONT_JP, font_size=Pt(22),
        font_color=BRAND["text_white"],
        bold=True,
    )
    # Right column body
    add_textbox(
        slide,
        left=Inches(7.2), top=Inches(3.1),
        width=Inches(4.8), height=Inches(3),
        text="右カラムの本文テキスト。比較や対照の情報を記載します。",
        font_name=FONT_JP, font_size=Pt(16),
        font_color=BRAND["neutral_200"],
    )


def create_code_slide(prs: Presentation):
    """スライド4: コードブロックスライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND["neutral_950"])

    # Top gradient bar
    add_gradient_line(slide, top=Pt(0), height=Pt(3))

    # Slide title
    add_textbox(
        slide,
        left=Inches(0.8), top=Inches(0.6),
        width=Inches(11), height=Inches(0.9),
        text="コードサンプル",
        font_name=FONT_JP, font_size=Pt(32),
        font_color=BRAND["text_white"],
        bold=True,
    )

    # Code block background
    code_bg = add_rounded_rect(
        slide,
        Inches(0.8), Inches(1.8),
        Inches(11.7), Inches(4.5),
        BRAND["neutral_700"],
    )

    # Terminal dots (red, yellow, green)
    dot_colors = ["#EF4444", "#F59E0B", "#10B981"]
    for i, color in enumerate(dot_colors):
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1.2) + Inches(0.35) * i, Inches(2.1),
            Pt(10), Pt(10),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = hex_to_rgb(color)
        dot.line.fill.background()

    # Code text
    code_box = add_textbox(
        slide,
        left=Inches(1.2), top=Inches(2.6),
        width=Inches(10.8), height=Inches(3.4),
        text="",
        font_name=FONT_MONO, font_size=Pt(14),
        font_color=BRAND["text_white"],
    )
    tf = code_box.text_frame
    tf.word_wrap = True

    code_lines = [
        ('// Vibe Coding with AI', BRAND["neutral_400"]),
        ('const vibeSession = await ai.pair({', BRAND["text_white"]),
        ('  model: "claude-opus-4-6",', BRAND["secondary"]),
        ('  mode: "creative",', BRAND["secondary"]),
        ('  live: true,', BRAND["accent"]),
        ('});', BRAND["text_white"]),
        ('', BRAND["text_white"]),
        ('console.log("Let\'s vibe! 🎵");', BRAND["primary"]),
    ]

    for i, (line, color) in enumerate(code_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.name = FONT_MONO
        p.font.size = Pt(15)
        p.font.color.rgb = hex_to_rgb(color)
        p.space_after = Pt(2)

    # Description below code
    add_textbox(
        slide,
        left=Inches(0.8), top=Inches(6.5),
        width=Inches(11), height=Inches(0.6),
        text="コードの説明テキストをここに記載",
        font_name=FONT_JP, font_size=Pt(14),
        font_color=BRAND["text_muted"],
    )


def create_image_text_slide(prs: Presentation):
    """スライド5: 画像 + テキストスライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND["neutral_950"])

    # Top gradient bar
    add_gradient_line(slide, top=Pt(0), height=Pt(3))

    # Slide title
    add_textbox(
        slide,
        left=Inches(0.8), top=Inches(0.6),
        width=Inches(11), height=Inches(0.9),
        text="画像 + テキスト",
        font_name=FONT_JP, font_size=Pt(32),
        font_color=BRAND["text_white"],
        bold=True,
    )

    # Divider
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(1.5),
        Inches(2), Pt(2),
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = hex_to_rgb(BRAND["primary"])
    divider.line.fill.background()

    # Left: image placeholder
    add_placeholder_image(
        slide,
        Inches(0.8), Inches(2.0),
        Inches(6), Inches(4.5),
    )

    # Right: text content
    add_textbox(
        slide,
        left=Inches(7.4), top=Inches(2.0),
        width=Inches(5.2), height=Inches(0.8),
        text="見出しテキスト",
        font_name=FONT_JP, font_size=Pt(24),
        font_color=BRAND["text_white"],
        bold=True,
    )

    add_textbox(
        slide,
        left=Inches(7.4), top=Inches(3.0),
        width=Inches(5.2), height=Inches(3.5),
        text="画像の説明や補足テキストをここに入力します。図表やスクリーンショットと一緒に使うと効果的です。",
        font_name=FONT_JP, font_size=Pt(16),
        font_color=BRAND["neutral_200"],
    )


def create_ending_slide(prs: Presentation):
    """スライド6: エンディングスライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND["neutral_950"])

    # Subtle glow
    glow = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(3), Inches(0.5),
        Inches(7), Inches(7),
    )
    glow.fill.solid()
    glow.fill.fore_color.rgb = hex_to_rgb(BRAND["secondary"])
    glow.fill.fore_color.brightness = 0.85
    glow.line.fill.background()
    from pptx.oxml.ns import qn
    fill_elem = glow._element.find(qn("p:spPr"))
    if fill_elem is not None:
        solid = fill_elem.find(qn("a:solidFill"))
        if solid is not None:
            srgb = solid.find(qn("a:srgbClr"))
            if srgb is not None:
                alpha = srgb.makeelement(qn("a:alpha"), {"val": "6000"})
                srgb.append(alpha)

    # Gradient bar at top
    add_gradient_line(slide, top=Pt(0), height=Pt(4))

    # Thank you text
    add_textbox(
        slide,
        left=Inches(0), top=Inches(2.2),
        width=SLIDE_WIDTH, height=Inches(1.5),
        text="Thank You",
        font_name=FONT_DISPLAY, font_size=Pt(56),
        font_color=BRAND["text_white"],
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    # Subtitle
    add_textbox(
        slide,
        left=Inches(0), top=Inches(3.8),
        width=SLIDE_WIDTH, height=Inches(0.8),
        text="AIと共に、コードで遊べ。",
        font_name=FONT_JP, font_size=Pt(22),
        font_color=BRAND["text_muted"],
        alignment=PP_ALIGN.CENTER,
    )

    # Contact / links
    add_textbox(
        slide,
        left=Inches(0), top=Inches(5.2),
        width=SLIDE_WIDTH, height=Inches(0.5),
        text="Live Vibe Coding Club",
        font_name=FONT_DISPLAY, font_size=Pt(14),
        font_color=BRAND["neutral_400"],
        alignment=PP_ALIGN.CENTER,
    )

    # Gradient bar at bottom
    add_gradient_line(slide, top=SLIDE_HEIGHT - Pt(4), height=Pt(4))


# ============================================================
# Main
# ============================================================

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Generate all slide types
    create_title_slide(prs)
    create_section_slide(prs)
    create_content_slide(prs)
    create_two_column_slide(prs)
    create_code_slide(prs)
    create_image_text_slide(prs)
    create_ending_slide(prs)

    output_path = "/home/kosawa/live-vibe-coding/brand/lvc-template.pptx"
    prs.save(output_path)
    print(f"PPTテンプレート生成完了: {output_path}")
    print(f"スライド数: {len(prs.slides)}")
    print()
    print("含まれるレイアウト:")
    print("  0. タイトルスライド")
    print("  1. セクション区切り")
    print("  2. コンテンツ（タイトル + 本文 + 箇条書き）")
    print("  3. 2カラム レイアウト")
    print("  4. コードブロック")
    print("  5. 画像 + テキスト")
    print("  6. エンディング")
    print()
    print("フォント要件:")
    print(f"  Display/Body: {FONT_DISPLAY}")
    print(f"  日本語: {FONT_JP}")
    print(f"  等幅: {FONT_MONO}")
    print("  ※ Google Fontsから事前インストールしてください")


if __name__ == "__main__":
    main()
